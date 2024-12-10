import google.generativeai as genai
import os
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation  # Import for generating PPT
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import textwrap

# Configure the API keys for both tasks
genai.configure(api_key="AIzaSyByWjZ1kOXvr_4vZtSxNQW_Rg-vVKmXNo0")  # First API key for basic content
genai.configure(api_key="AIzaSyDENy-P2AkLAy0CGi6Zx14Y5H80-yNFkfE")  # Second API key for detailed descriptions


def generate_ppt_from_prompt(prompt):
    """Generate a PowerPoint presentation using content from Gemini AI based on a user-provided prompt."""
    try:
        # Ensure the prompt is not empty
        if not prompt.strip():
            messagebox.showerror("Error", "Prompt cannot be empty. Please enter a valid prompt.")
            return

        # Generate content from the prompt using Gemini AI (First API Key)
        model = genai.GenerativeModel(model_name="gemini-1.5-pro")
        response = model.generate_content([prompt])
        content = response.text

        if not content:
            content = "No content was generated. Please try again with a different prompt."

        # Use the second API key to generate detailed descriptions
        detailed_descriptions = get_detailed_descriptions(content)

        # Create a PowerPoint presentation
        prs = Presentation()

        # Set the theme for the presentation
        set_ppt_theme(prs)

        # Split the content into sections based on paragraphs
        sections = [sec.strip() for sec in content.split("\n\n") if sec.strip()]

        # Create slides based on content
        for i, section in enumerate(sections):
            if i == 0:  # Title Slide
                add_title_slide(prs, section)
            else:  # Content slides
                description = detailed_descriptions[i - 1] if i - 1 < len(
                    detailed_descriptions) else "No detailed description."
                add_content_slide(prs, section, description)

        # Save PPT file
        ppt_filename = filedialog.asksaveasfilename(defaultextension=".pptx",
                                                    filetypes=[("PowerPoint Files", "*.pptx")])
        if ppt_filename:
            prs.save(ppt_filename)
            messagebox.showinfo("Success", f"PowerPoint presentation saved as {ppt_filename}")
    except Exception as e:
        messagebox.showerror("Error", f"Failed to generate PowerPoint: {str(e)}")


def set_ppt_theme(prs):
    """Set the theme and design for the PowerPoint presentation."""
    # Create a consistent background color for the presentation
    for slide in prs.slides:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray background


def add_title_slide(prs, content):
    """Create the title slide."""
    slide_layout = prs.slide_layouts[0]  # Title Slide Layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Generated Presentation"

    # Set subtitle (content) for the title slide
    subtitle = slide.placeholders[1]
    subtitle.text = content


def add_content_slide(prs, content, description):
    """Create a content slide with a header and the detailed description."""
    slide_layout = prs.slide_layouts[1]  # Title and Content Layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Detailed Content"

    # Create a text box with a max height to prevent overflow
    textbox = slide.shapes.placeholders[1]
    text_frame = textbox.text_frame
    text_frame.clear()

    # Wrap the content to fit within the slide width
    wrapped_content = textwrap.fill(content, width=80)  # Wrap text to 80 characters per line
    wrapped_description = textwrap.fill(description, width=80)  # Wrap description text

    # Add wrapped text to the slide
    p = text_frame.add_paragraph()
    p.text = f"{wrapped_content}\n\nDetailed Description:\n{wrapped_description}"
    p.font.size = Pt(18)
    p.font.name = "Calibri"

    # Adjust the font size if the text is too large to fit
    adjust_font_size(text_frame)


def adjust_font_size(text_frame):
    """Adjust the font size to fit the text within the slide."""
    max_font_size = Pt(18)  # Maximum font size
    min_font_size = Pt(10)  # Minimum font size

    # Count the number of lines in the text_frame and adjust font size if too many lines
    line_count = len(text_frame.text.split("\n"))

    # Reduce font size based on the number of lines
    if line_count > 10:
        max_font_size = Pt(14)  # Reduce font size if the text has too many lines
    elif line_count > 5:
        max_font_size = Pt(16)  # Reduce slightly for moderate amount of text

    # Apply the font size adjustment
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = max_font_size


def get_detailed_descriptions(content):
    """Use the second API key to generate detailed descriptions based on content."""
    try:
        # Send request to the second API key to generate detailed descriptions for each section
        # For now, I am assuming a simple case where we split the content and generate descriptions for each section.

        detailed_descriptions = []
        sections = content.split("\n\n")

        for section in sections:
            response = genai.GenerativeModel(model_name="gemini-1.5-pro").generate_content(
                [f"Provide a detailed description for: {section}"])
            detailed_descriptions.append(response.text)

        return detailed_descriptions

    except Exception as e:
        print(f"Error while generating detailed descriptions: {e}")
        return []


# Create the GUI
def run_gui():
    root = tk.Tk()
    root.title("Gemini AI PPT Generator")
    root.geometry("800x600")
    root.configure(bg="#f4f4f9")

    # Main Frame for Layout
    main_frame = tk.Frame(root, bg="#ffffff", padx=20, pady=20)
    main_frame.pack(fill=tk.BOTH, expand=True)

    # Title Section
    title_frame = tk.Frame(main_frame, bg="#4c8bf5", bd=0)
    title_frame.grid(row=0, column=0, columnspan=2, pady=10)

    title_label = tk.Label(
        title_frame,
        text="Gemini AI PPT Generator",
        font=("Helvetica", 24, "bold"),
        bg="#4c8bf5",
        fg="white",
        padx=20,
        pady=20,
    )
    title_label.pack()

    # Prompt input
    prompt_label = tk.Label(main_frame, text="Enter a Prompt for PowerPoint Generation:", font=("Helvetica", 12),
                            bg="#f4f4f9")
    prompt_label.grid(row=1, column=0, pady=10, sticky="w")

    prompt_entry = tk.Entry(main_frame, font=("Helvetica", 12), width=60)
    prompt_entry.grid(row=1, column=1, pady=10)

    # Button to generate PowerPoint
    generate_ppt_button = tk.Button(
        main_frame,
        text="Generate PowerPoint from Prompt",
        font=("Helvetica", 14),
        bg="#3e64ff",
        fg="white",
        padx=20,
        pady=10,
        command=lambda: generate_ppt_from_prompt(prompt_entry.get())
    )
    generate_ppt_button.grid(row=2, column=0, columnspan=2, pady=20)

    # Footer Section
    footer_label = tk.Label(
        root,
        text="Powered by Gemini AI | Developed by You",
        font=("Helvetica", 10, "italic"),
        bg="#4c8bf5",
        fg="white",
        pady=5,
    )
    footer_label.pack(side=tk.BOTTOM, fill=tk.X)

    root.mainloop()


if __name__ == "__main__":
    run_gui()
